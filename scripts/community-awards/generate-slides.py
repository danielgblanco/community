import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.enum.text import PP_ALIGN
import requests
from io import BytesIO
import pandas as pd
import sys
from PIL import Image
import random

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# URL to the OpenTelemetry logo
OPENTELEMETRY_LOGO_URL = 'https://opentelemetry.io/img/social/logo-wordmark-001.png'


def download_image(url):
    """Downloads an image from a URL and returns it as a BytesIO object."""
    try:
        response = requests.get(url)
        response.raise_for_status()
        logging.info(f"Successfully downloaded image from {url}")
        return BytesIO(response.content)
    except requests.RequestException:
        logging.error(f"Failed to download image from {url}")
        return None


def get_github_profile_picture(username):
    """
    Fetches the GitHub profile picture for a given username.
    Returns an image in BytesIO format if successful, otherwise None.
    """
    url = f"https://github.com/{username}.png"
    return download_image(url)


def add_fade_transition(slide):
    """Adds a fade-in and fade-out transition to the slide."""
    transition_xml = (
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'spd="slow" advTm="3000" dur="1000">'
        '<p:fade/></p:transition>'
    )
    slide._element.append(parse_xml(transition_xml))


def set_slide_duration(slide, duration_sec):
    """Sets the duration for which the slide is displayed."""
    timing_xml = (
        '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        f'<p:tnLst><p:par><p:cTn dur="{int(duration_sec * 1000)}" restart="never" nodeType="tmRoot"/></p:par></p:tnLst>'
        '</p:timing>'
    )
    slide._element.append(parse_xml(timing_xml))


def add_footer(slide):
    """Adds a footer to the slide."""
    footer_text = "OpenTelemetry Community Awards"
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = footer_text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)  # Grey color


def create_presentation(feedback_data, output_file):
    """Creates a PowerPoint presentation with feedback slides for each user."""
    prs = Presentation()

    # Download the OpenTelemetry logo
    logo_image = download_image(OPENTELEMETRY_LOGO_URL)
    if not logo_image:
        logging.error("OpenTelemetry logo could not be downloaded. Exiting.")
        sys.exit(1)

    # Get the aspect ratio of the logo
    logo_image.seek(0)
    logo = Image.open(logo_image)
    logo_aspect_ratio = logo.width / logo.height

    # Shuffle the feedback data
    feedback_data = feedback_data.sample(frac=1).reset_index(drop=True)

    for _, row in feedback_data.iterrows():
        username = row.iloc[0]
        feedback = row.iloc[1]

        # Add a new slide with a title layout
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Add a textbox for feedback content with soft gray text color and adjusted font size
        text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3.5))  # Adjusted height
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.add_paragraph()
        p.text = feedback
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(80, 80, 80)  # Soft gray color for readability

        # Try to add the GitHub profile picture if available
        profile_pic = get_github_profile_picture(username)
        if profile_pic:
            profile_pic_height = Inches(1.5)
            profile_pic_width = int(profile_pic_height *
                                    (Image.open(profile_pic).width / Image.open(profile_pic).height))
            slide.shapes.add_picture(profile_pic, Inches(0.5), Inches(0.5), profile_pic_width, profile_pic_height)

            # Move the title below the profile picture
            title = slide.shapes.title
            title.left = Inches(0.5)
            title.top = Inches(2.2)
            title.width = Inches(3)  # Set a fixed width for the title
            title.height = Inches(0.5)
            title.text = username
            title.text_frame.paragraphs[0].font.size = Pt(24)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)  # Soft blue color
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Align text to the left

            # Add the OpenTelemetry logo to the same vertical position but on the opposite side
            slide_width = prs.slide_width
            logo_height = profile_pic_height
            logo_width = logo_height * logo_aspect_ratio
            slide.shapes.add_picture(logo_image, slide_width - logo_width - Inches(0.5), Inches(0.5),
                                     logo_width, logo_height)

        # Add fade-in and fade-out transition
        add_fade_transition(slide)

        # Set slide duration to 3 seconds
        set_slide_duration(slide, 15)

        # Add footer
        add_footer(slide)

        logging.info(f"Added slide for {username}")

    # Save the presentation
    prs.save(output_file)
    logging.info(f"Presentation created: {output_file}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python generate-slides.py <input_csv_file> <output_pptx_file>")
        sys.exit(1)

    input_csv_file = sys.argv[1]
    output_pptx_file = sys.argv[2]

    # Load feedback data from the specified CSV file
    feedback_data = pd.read_csv(input_csv_file, header=0, usecols=[0, 1])

    create_presentation(feedback_data, output_pptx_file)
